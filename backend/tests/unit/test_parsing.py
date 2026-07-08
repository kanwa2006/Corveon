"""Unit tests for PDF parsing (app/ingestion/parsing.py). Test PDFs are
generated in-memory with PyMuPDF itself rather than checked-in binary
fixtures — genuinely parseable, no fixture-drift risk."""

from __future__ import annotations

import io

import fitz
import pytest
from app.ingestion.parsing import (
    MAX_PAGE_COUNT,
    MAX_PARAGRAPH_COUNT,
    MAX_SLIDE_COUNT,
    DocumentParseError,
    DocumentTooLargeError,
    UnsupportedDocumentTypeError,
    parse_document,
    parse_docx,
    parse_image,
    parse_markdown,
    parse_pdf,
    parse_pptx,
)
from docx import Document as DocxDocument
from PIL import Image, ImageDraw
from pptx import Presentation

pytestmark = pytest.mark.unit


def _make_pdf(page_texts: list[str], *, encrypt: bool = False) -> bytes:
    doc = fitz.open()
    for text in page_texts:
        page = doc.new_page()
        page.insert_text((72, 72), text)
    data: bytes = (
        doc.tobytes(
            encryption=fitz.PDF_ENCRYPT_AES_256,
            owner_pw="owner-secret",
            user_pw="user-secret",  # a user password is what actually gates opening the file
        )
        if encrypt
        else doc.tobytes()
    )
    doc.close()
    return data


def test_parse_pdf_extracts_text_and_page_count() -> None:
    pdf_bytes = _make_pdf(["First page text.", "Second page text."])
    parsed = parse_pdf(pdf_bytes)
    assert parsed.page_count == 2
    assert "First page text." in parsed.pages[0]
    assert "Second page text." in parsed.pages[1]


def test_parse_pdf_rejects_non_pdf_bytes() -> None:
    with pytest.raises(DocumentParseError):
        parse_pdf(b"this is not a pdf at all")


def test_parse_pdf_rejects_encrypted_pdf() -> None:
    encrypted = _make_pdf(["secret"], encrypt=True)
    with pytest.raises(DocumentParseError, match="password-protected"):
        parse_pdf(encrypted)


def test_parse_pdf_rejects_documents_over_page_cap() -> None:
    pdf_bytes = _make_pdf(["page"] * (MAX_PAGE_COUNT + 1))
    with pytest.raises(DocumentTooLargeError):
        parse_pdf(pdf_bytes)


def _text_image_bytes(text: str, *, size: tuple[int, int] = (600, 200)) -> bytes:
    image = Image.new("RGB", size, color="white")
    draw = ImageDraw.Draw(image)
    draw.text((20, size[1] // 2 - 20), text, fill="black")
    buf = io.BytesIO()
    image.save(buf, format="PNG")
    return buf.getvalue()


def test_parse_pdf_ocrs_a_scanned_page_with_no_text_layer() -> None:
    doc = fitz.open()
    page = doc.new_page()
    page.insert_image(fitz.Rect(0, 0, 600, 200), stream=_text_image_bytes("SCANNED PAGE"))
    pdf_bytes = doc.tobytes()
    doc.close()

    parsed = parse_pdf(pdf_bytes)
    assert parsed.page_count == 1
    assert "SCANNED" in parsed.pages[0].upper()


# ── DOCX ─────────────────────────────────────────────────────────────────


def _make_docx(paragraphs: list[str]) -> bytes:
    doc = DocxDocument()
    for text in paragraphs:
        doc.add_paragraph(text)
    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


def test_parse_docx_extracts_paragraph_text() -> None:
    docx_bytes = _make_docx(["Metformin is a first-line treatment.", "Second paragraph."])
    parsed = parse_docx(docx_bytes)
    assert parsed.page_count == 1
    assert "Metformin is a first-line treatment." in parsed.pages[0]
    assert "Second paragraph." in parsed.pages[0]


def test_parse_docx_rejects_non_docx_bytes() -> None:
    with pytest.raises(DocumentParseError):
        parse_docx(b"this is not a docx at all")


def test_parse_docx_rejects_documents_over_paragraph_cap() -> None:
    docx_bytes = _make_docx(["paragraph"] * (MAX_PARAGRAPH_COUNT + 1))
    with pytest.raises(DocumentTooLargeError):
        parse_docx(docx_bytes)


# ── PPTX ─────────────────────────────────────────────────────────────────


def _make_pptx(slide_titles: list[str]) -> bytes:
    presentation = Presentation()
    layout = presentation.slide_layouts[1]
    for title in slide_titles:
        slide = presentation.slides.add_slide(layout)
        slide.shapes.title.text = title
    buf = io.BytesIO()
    presentation.save(buf)
    return buf.getvalue()


def test_parse_pptx_extracts_one_page_per_slide() -> None:
    pptx_bytes = _make_pptx(["Renal dosing", "Beers criteria"])
    parsed = parse_pptx(pptx_bytes)
    assert parsed.page_count == 2
    assert "Renal dosing" in parsed.pages[0]
    assert "Beers criteria" in parsed.pages[1]


def test_parse_pptx_rejects_non_pptx_bytes() -> None:
    with pytest.raises(DocumentParseError):
        parse_pptx(b"this is not a pptx at all")


def test_parse_pptx_rejects_presentations_over_slide_cap() -> None:
    pptx_bytes = _make_pptx(["slide"] * (MAX_SLIDE_COUNT + 1))
    with pytest.raises(DocumentTooLargeError):
        parse_pptx(pptx_bytes)


# ── Markdown ─────────────────────────────────────────────────────────────


def test_parse_markdown_preserves_raw_text() -> None:
    parsed = parse_markdown(b"# Renal Dosing\n\nMetformin is contraindicated below eGFR 30.")
    assert parsed.page_count == 1
    assert parsed.pages[0] == "# Renal Dosing\n\nMetformin is contraindicated below eGFR 30."


def test_parse_markdown_rejects_non_utf8_bytes() -> None:
    with pytest.raises(DocumentParseError):
        parse_markdown(b"\xff\xfe\x00\x81invalid")


# ── Images ───────────────────────────────────────────────────────────────


def test_parse_image_ocrs_text_from_a_png() -> None:
    parsed = parse_image(_text_image_bytes("HELLO"))
    assert parsed.page_count == 1
    assert "HELLO" in parsed.pages[0].upper()


def test_parse_image_rejects_non_image_bytes() -> None:
    with pytest.raises(DocumentParseError):
        parse_image(b"this is not an image at all")


# ── Registry dispatch ────────────────────────────────────────────────────


def test_parse_document_dispatches_by_mime_type() -> None:
    docx_bytes = _make_docx(["hello"])
    parsed = parse_document(
        docx_bytes, "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
    )
    assert "hello" in parsed.pages[0]


def test_parse_document_raises_for_unregistered_mime_type() -> None:
    with pytest.raises(UnsupportedDocumentTypeError):
        parse_document(b"data", "application/x-not-a-real-format")
