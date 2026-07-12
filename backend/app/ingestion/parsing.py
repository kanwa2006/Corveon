"""Multi-format document parsing (docs/ARCHITECTURE.md §3, §9; CLAUDE.md
Month 1 roadmap: "Multi-format ingestion (DOCX/PPT/MD/images + OCR)"). Every
parser returns the same ``ParsedDocument`` shape (page_count + a list of
per-page/per-unit text) so ``chunk_pages`` (app/ingestion/chunking.py) and
everything downstream stays format-agnostic. ``parse_document`` is the single
entry point, dispatching on MIME type — the parser registry the ingestion
worker and the upload endpoint both go through, so adding a format later
means adding one dispatch entry, not touching either caller."""

from __future__ import annotations

import io
import zipfile
from dataclasses import dataclass

import fitz  # PyMuPDF
import pytesseract
from docx import Document as DocxDocument
from docx.opc.exceptions import PackageNotFoundError as DocxPackageNotFoundError
from PIL import Image, UnidentifiedImageError
from pptx import Presentation
from pptx.exc import PackageNotFoundError as PptxPackageNotFoundError

# PDF-bomb defense (docs/SECURITY.md: "image/PDF-bomb limits"). A legitimate
# clinical document is not thousands of pages; this bounds worst-case parse
# time/memory for a maliciously crafted upload. The same reasoning sets the
# per-format caps below for DOCX/PPTX/images.
MAX_PAGE_COUNT = 500
MAX_PARAGRAPH_COUNT = 20_000
MAX_SLIDE_COUNT = 500
MAX_IMAGE_PIXELS = 30_000_000  # ~6000x5000 — generous for a scanned page, not for a bomb
# 2x zoom balances OCR accuracy against per-page render/OCR cost; PyMuPDF's
# default (72 DPI) renders text too small for Tesseract to read reliably.
_OCR_ZOOM = 2.0

_DOCX_MIME = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
_PPTX_MIME = "application/vnd.openxmlformats-officedocument.presentationml.presentation"


class DocumentParseError(Exception):
    """The file could not be parsed (corrupt/encrypted/not a valid file of
    its declared type)."""


class DocumentTooLargeError(Exception):
    """The document exceeds a format's page/paragraph/slide/pixel cap — a
    bomb defense."""


class UnsupportedDocumentTypeError(Exception):
    """No parser is registered for this MIME type."""


@dataclass(frozen=True, slots=True)
class ParsedDocument:
    page_count: int
    pages: list[str]


def _ocr_image(image: Image.Image) -> str:
    if image.width * image.height > MAX_IMAGE_PIXELS:
        raise DocumentTooLargeError(
            f"Image is {image.width}x{image.height} ({image.width * image.height:,} px), "
            f"exceeding the {MAX_IMAGE_PIXELS:,}px limit."
        )
    return pytesseract.image_to_string(image)  # type: ignore[no-any-return]


def _ocr_pdf_page(page: fitz.Page) -> str:
    """Renders a page with no extractable text layer to an image and OCRs
    it — the scanned-document fallback. Never runs on a page that already
    has a text layer (checked by the caller), so a normal text PDF pays
    zero OCR cost."""
    # Enforce the pixel cap BEFORE rasterizing: get_pixmap allocates the
    # full bitmap, so a small PDF declaring an extreme MediaBox would OOM
    # the worker before the post-render check in _ocr_image ever ran.
    rendered_pixels = page.rect.width * _OCR_ZOOM * page.rect.height * _OCR_ZOOM
    if rendered_pixels > MAX_IMAGE_PIXELS:
        raise DocumentTooLargeError(
            f"PDF page is {page.rect.width:.0f}x{page.rect.height:.0f} pt "
            f"(~{rendered_pixels:,.0f} px at OCR resolution), exceeding the "
            f"{MAX_IMAGE_PIXELS:,}px limit."
        )
    pixmap = page.get_pixmap(matrix=fitz.Matrix(_OCR_ZOOM, _OCR_ZOOM))
    image = Image.open(io.BytesIO(pixmap.tobytes("png")))
    return _ocr_image(image)


def parse_pdf(data: bytes) -> ParsedDocument:
    try:
        doc = fitz.open(stream=data, filetype="pdf")
    except Exception as exc:  # PyMuPDF's own error types vary by failure mode
        raise DocumentParseError(f"Could not open PDF: {exc}") from exc

    try:
        if doc.is_encrypted:
            raise DocumentParseError("PDF is password-protected.")
        if doc.page_count > MAX_PAGE_COUNT:
            raise DocumentTooLargeError(
                f"PDF has {doc.page_count} pages, exceeding the {MAX_PAGE_COUNT}-page limit."
            )
        pages = []
        for page in doc:
            text = page.get_text()
            # A page with no extractable text is a scanned image, not an
            # empty page (docs rarely have genuinely blank pages worth
            # uploading) — OCR it rather than silently dropping its content.
            pages.append(text if text.strip() else _ocr_pdf_page(page))
        return ParsedDocument(page_count=doc.page_count, pages=pages)
    finally:
        doc.close()


def parse_docx(data: bytes) -> ParsedDocument:
    try:
        doc = DocxDocument(io.BytesIO(data))
    except (DocxPackageNotFoundError, zipfile.BadZipFile, KeyError, ValueError) as exc:
        raise DocumentParseError(f"Could not open DOCX: {exc}") from exc

    if len(doc.paragraphs) > MAX_PARAGRAPH_COUNT:
        raise DocumentTooLargeError(
            f"DOCX has {len(doc.paragraphs)} paragraphs, exceeding the {MAX_PARAGRAPH_COUNT} limit."
        )
    # Word has no fixed page concept (text reflows) — the whole document is
    # one logical unit; chunk_pages' own paragraph-aware splitting handles
    # breaking it up from here.
    text = "\n\n".join(p.text for p in doc.paragraphs if p.text.strip())
    return ParsedDocument(page_count=1, pages=[text])


def parse_pptx(data: bytes) -> ParsedDocument:
    try:
        presentation = Presentation(io.BytesIO(data))
    except (PptxPackageNotFoundError, zipfile.BadZipFile, KeyError, ValueError) as exc:
        raise DocumentParseError(f"Could not open PPTX: {exc}") from exc

    slides = list(presentation.slides)
    if len(slides) > MAX_SLIDE_COUNT:
        raise DocumentTooLargeError(
            f"PPTX has {len(slides)} slides, exceeding the {MAX_SLIDE_COUNT}-slide limit."
        )
    pages = []
    for slide in slides:
        texts = [
            shape.text_frame.text
            for shape in slide.shapes
            if shape.has_text_frame and shape.text_frame.text.strip()
        ]
        pages.append("\n\n".join(texts))
    return ParsedDocument(page_count=len(slides), pages=pages)


def parse_markdown(data: bytes) -> ParsedDocument:
    try:
        text = data.decode("utf-8")
    except UnicodeDecodeError as exc:
        raise DocumentParseError(f"Could not decode Markdown as UTF-8: {exc}") from exc
    # Kept as raw Markdown, not rendered to HTML/plain text — headings and
    # structure are useful signal for both chunking and the LLM itself,
    # which reads Markdown natively.
    return ParsedDocument(page_count=1, pages=[text])


def parse_image(data: bytes) -> ParsedDocument:
    try:
        image = Image.open(io.BytesIO(data))
        image.load()  # Pillow defers decoding; force it now to catch truncated/bad files here
    except (UnidentifiedImageError, OSError) as exc:
        raise DocumentParseError(f"Could not open image: {exc}") from exc
    return ParsedDocument(page_count=1, pages=[_ocr_image(image)])


_PARSERS = {
    "application/pdf": parse_pdf,
    _DOCX_MIME: parse_docx,
    _PPTX_MIME: parse_pptx,
    "text/markdown": parse_markdown,
    "text/x-markdown": parse_markdown,
    "image/png": parse_image,
    "image/jpeg": parse_image,
}


def parse_document(data: bytes, mime_type: str) -> ParsedDocument:
    """The parser registry — dispatches on MIME type. Raises
    UnsupportedDocumentTypeError for anything not registered; callers should
    validate mime_type against the same allow-list at upload time rather
    than relying on this raising (docs/SECURITY.md: validate at the
    boundary), but this raises rather than silently mis-parsing regardless."""
    parser = _PARSERS.get(mime_type)
    if parser is None:
        raise UnsupportedDocumentTypeError(f"No parser registered for MIME type {mime_type!r}.")
    return parser(data)
